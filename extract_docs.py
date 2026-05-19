import os
import glob
from docx import Document
from pptx import Presentation
import pypdfium2 as pdfium
import openpyxl

def extract_docx(filepath):
    try:
        doc = Document(filepath)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip() != ""])
    except Exception as e:
        return f"Error reading {filepath}: {e}"

def extract_pptx(filepath):
    try:
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"Error reading {filepath}: {e}"

def extract_pdf(filepath):
    try:
        pdf = pdfium.PdfDocument(filepath)
        text = []
        for page in pdf:
            textpage = page.get_textpage()
            text.append(textpage.get_text_bounded())
        return "\n".join(text)
    except Exception as e:
        return f"Error reading {filepath}: {e}"

def extract_xlsx(filepath):
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) for cell in row if cell is not None]
                if row_text:
                    text.append(" | ".join(row_text))
        return "\n".join(text)
    except Exception as e:
        return f"Error reading {filepath}: {e}"

def extract_all(directory):
    output_dir = os.path.join(directory, "extracted_texts")
    os.makedirs(output_dir, exist_ok=True)
    
    files = []
    for ext in ("*.docx", "*.pptx", "*.pdf", "*.xlsx"):
        files.extend(glob.glob(os.path.join(directory, "**", ext), recursive=True))
        
    for filepath in files:
        if "extracted_texts" in filepath:
            continue
            
        print(f"Extracting {os.path.basename(filepath)}...")
        ext = os.path.splitext(filepath)[1].lower()
        content = ""
        
        if ext == ".docx":
            content = extract_docx(filepath)
        elif ext == ".pptx":
            content = extract_pptx(filepath)
        elif ext == ".pdf":
            content = extract_pdf(filepath)
        elif ext == ".xlsx":
            content = extract_xlsx(filepath)
            
        if content.strip():
            out_name = os.path.basename(filepath) + ".txt"
            out_path = os.path.join(output_dir, out_name)
            with open(out_path, "w", encoding="utf-8") as f:
                f.write(content)

if __name__ == "__main__":
    extract_all("/Users/jakubbednar/Documents/ivp-upv/škola")
