from pptx import Presentation
import copy
import os

print("CWD is:", os.getcwd())

prs = Presentation("CZU_IVP_prezentace_sablona_vzor.pptx")

def replace_title_runs(slide, shape_name, new_paragraphs):
    for shape in slide.placeholders:
        if shape.name == shape_name:
            for p_idx, text in enumerate(new_paragraphs):
                if p_idx < len(shape.text_frame.paragraphs):
                    p = shape.text_frame.paragraphs[p_idx]
                    if p.runs:
                        p.runs[0].text = text
                        for r in p.runs[1:]:
                            r.text = ""
                    else:
                        p.text = text

def set_bullets(slide, shape_idx, bullet_lines):
    target_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == shape_idx:
            target_shape = shape
            break
            
    if not target_shape or not target_shape.has_text_frame:
        return
        
    tf = target_shape.text_frame
    
    if not tf.paragraphs:
        return
        
    p0 = tf.paragraphs[0]
    
    if bullet_lines:
        p0.text = bullet_lines[0]
        
    for line in bullet_lines[1:]:
        new_p = tf.add_paragraph()
        new_p.text = line
        new_p.level = p0.level
        if p0._p.pPr is not None:
            new_p._p.insert(0, copy.deepcopy(p0._p.pPr))

def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

vysledky_layout = prs.slides[4].slide_layout
slide_v2 = prs.slides.add_slide(vysledky_layout)
slide_v3 = prs.slides.add_slide(vysledky_layout)
slide_v4 = prs.slides.add_slide(vysledky_layout)

move_slide(prs, 8, 5)
move_slide(prs, 9, 6)
move_slide(prs, 10, 7)

# Slide 0: Title
replace_title_runs(prs.slides[0], "Zástupný text 2", ["Využití platformy umělé inteligence ve výuce odborných předmětů", "Jakub Bednář"])
replace_title_runs(prs.slides[0], "Zástupný text 3", ["Vedoucí práce: [Doplň jméno]", "Oponent práce: [Doplň jméno]"])

# Slide 1: Obsah
set_bullets(prs.slides[1], 15, [
    "1. Úvod a cíle práce",
    "2. Metodika výzkumu",
    "3. Hlavní výsledky a paradoxy",
    "4. Praktická doporučení a závěr"
])

# Slide 2: Cíl
set_bullets(prs.slides[2], 15, [
    "Zmapovat reálný stav využívání AI (ChatGPT aj.) na středních odborných školách.",
    "Analyzovat postoje studentů vs. pedagogů k AI jako vzdělávacímu nástroji.",
    "Navrhnout efektivní začlenění AI pro podporu analytického myšlení, nikoliv jen jako usnadnění práce."
])

# Slide 3: Metodika
set_bullets(prs.slides[3], 15, [
    "Design smíšeného výzkumu (Mixed Methods).",
    "Kvantitativní část: Dotazníkové šetření mezi 240 studenty SOŠ (výrazné zastoupení technických oborů a IT).",
    "Kvalitativní část: Polostrukturované rozhovory s 12 pedagogy (různé aprobace a délky praxe)."
])

# Slide 4: Výsledky 1
for shape in prs.slides[4].placeholders:
    if shape.placeholder_format.idx == 10: shape.text = "Výsledky (1/4): Využívání AI"
set_bullets(prs.slides[4], 15, [
    "76,7 % studentů využívá AI alespoň jednou týdně.",
    "Primární využití: vysvětlování složité látky (AI jako interaktivní tutor).",
    "AI není primárně využívána k obcházení práce, ale k hlubšímu pochopení učiva."
])

# Slide 5: Výsledky 2
for shape in prs.slides[5].placeholders:
    if shape.placeholder_format.idx == 10: shape.text = "Výsledky (2/4): Kritický přístup"
set_bullets(prs.slides[5], 15, [
    "43 % studentů si aktivně ověřuje výstupy z AI pomocí jiných zdrojů.",
    "Studenti k AI přistupují s mnohem větší opatrností, než pedagogové předpokládají.",
    "Pozorujeme přirozený vývoj k analytickému využití technologie, nikoliv k pouhému přejímání chyb."
])

# Slide 6: Výsledky 3
for shape in prs.slides[6].placeholders:
    if shape.placeholder_format.idx == 10: shape.text = "Výsledky (3/4): Institucionální vakuum"
set_bullets(prs.slides[6], 15, [
    "Identifikován paradox „Institucionálního vakua“ na středních školách.",
    "Téměř 40 % studentů vůbec nezná pravidla své školy pro práci s AI.",
    "Školy nereagují dostatečně rychle, což vede k nejasnostem v tom, co je a není etické využití."
])

# Slide 7: Výsledky 4
for shape in prs.slides[7].placeholders:
    if shape.placeholder_format.idx == 10: shape.text = "Výsledky (4/4): Stres pedagogů"
set_bullets(prs.slides[7], 15, [
    "Pedagogové pociťují ztrátu kontroly nad klasifikací a prožívají z AI stres.",
    "Dva tábory učitelů: „striktní zakazovači“ vs. „proaktivní integrátoři“ (často z IT).",
    "Integrátoři mění systém: Těžištěm už není hotový produkt, ale formativní hodnocení a schopnost obhájit postup."
])

# Slide 8: Závěr a přínos
set_bullets(prs.slides[8], 15, [
    "Nutnost transformovat roli pedagoga: z „přenašeče informací“ na facilitátora učení.",
    "AI se ukazuje jako silný nástroj pro projektovou a kooperativní výuku.",
    "Školy musí nahradit plošné zákazy smysluplnou metodickou integrací.",
    "Závěr: „Kdo se s AI naučí pracovat, neztratí práci. Nahradí nás ti, kteří AI používají.“"
])

# Slide 9: Zdroje
set_bullets(prs.slides[9], 15, [
    "NPI ČR (2023). Doporučení pro využívání umělé inteligence ve školách.",
    "HOLMES, W. et al. (2022). Artificial Intelligence in Education.",
    "LUCKIN, R. et al. (2016). Intelligence Unleashed: An argument for AI in Education."
])

# Slide 10: Otázky k obhajobě
set_bullets(prs.slides[10], 15, [
    "Děkuji za pozornost.",
    "Zde je prostor pro otázky oponenta a vážené komise."
])

prs.save("BP_Prezentace_Final_Sablona_OK.pptx")
print("Hotovo")
